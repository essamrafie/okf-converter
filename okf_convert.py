#!/usr/bin/env python3
"""
okf_convert.py — Convert any directory into an OKF-conformant knowledge bundle.

Supported input formats:
  Text   : .md .txt .rst .py .js .ts .sh .yaml .yml .json .toml .cfg .ini .env
  Office : .docx .pptx .xlsx .xls .xlsm .csv .tsv
  PDF    : .pdf  (text-layer extraction; falls back gracefully for scanned PDFs)

Usage:
    python okf_convert.py --input ./my-docs --output ./okf-bundle
    python okf_convert.py --input ./my-docs --output ./okf-bundle \\
        --endpoint http://litellm.internal/v1 --model qwen3-30b
    python okf_convert.py --input ./my-docs --output ./okf-bundle --dry-run
"""

import argparse
import hashlib
import json
import os
import re
import sys
import textwrap
from datetime import datetime, timezone
from pathlib import Path

import requests
import yaml

# ─────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────

TEXT_EXTENSIONS = {
    ".md", ".txt", ".rst",
    ".py", ".js", ".ts", ".sh",
    ".yaml", ".yml", ".json", ".toml", ".cfg", ".ini", ".env",
}

OFFICE_EXTENSIONS = {
    ".docx", ".pptx",
    ".xlsx", ".xls", ".xlsm", ".csv", ".tsv",
}

PDF_EXTENSIONS = {".pdf"}

ALL_SUPPORTED = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | PDF_EXTENSIONS

SKIP_DIRS = {
    ".git", ".svn", "__pycache__", "node_modules", ".venv",
    "venv", "dist", "build", ".mypy_cache", ".pytest_cache",
}

RESERVED_FILENAMES = {"index", "log"}

MANIFEST_FILE     = ".okf_manifest.json"   # tracks source path → sha256 hash

MAX_CONTENT_CHARS = 6000
DEFAULT_ENDPOINT  = "http://localhost:4000/v1"
DEFAULT_MODEL     = "gpt-4o-mini"

EXTENSION_TYPE_HINT = {
    # Text
    ".md":   "Document",    ".txt":  "Document",    ".rst":  "Document",
    ".py":   "Code Module", ".js":   "Code Module",  ".ts":  "Code Module",
    ".sh":   "Runbook",
    ".yaml": "Configuration", ".yml": "Configuration", ".json": "Configuration",
    ".toml": "Configuration", ".cfg": "Configuration",
    ".ini":  "Configuration", ".env": "Configuration",
    # Office
    ".docx": "Document",    ".pptx": "Playbook",
    ".xlsx": "Dataset",     ".xls":  "Dataset",
    ".xlsm": "Dataset",     ".csv":  "Dataset",    ".tsv":  "Dataset",
    # PDF
    ".pdf":  "Document",
}

CODE_EXTENSIONS = {".py", ".js", ".ts", ".sh"}
CONFIG_EXTENSIONS = {".yaml", ".yml", ".json", ".toml", ".cfg", ".ini", ".env"}

# ─────────────────────────────────────────────────────────────
# Text extractors
# ─────────────────────────────────────────────────────────────

def extract_text(src: Path) -> tuple[str, str]:
    """
    Returns (extracted_text, extraction_note).
    extraction_note is a short label like 'docx', 'pdf-text', 'xlsx-3-sheets'
    to add context to the OKF body.
    """
    ext = src.suffix.lower()

    # ── Plain text / code / config ──────────────────────────
    if ext in TEXT_EXTENSIONS:
        try:
            return src.read_text(encoding="utf-8", errors="replace"), ext.lstrip(".")
        except Exception as e:
            return f"[read error: {e}]", "error"

    # ── DOCX ────────────────────────────────────────────────
    if ext == ".docx":
        try:
            import docx as python_docx
            doc = python_docx.Document(str(src))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else ""
                    if style.startswith("Heading"):
                        level = style.replace("Heading ", "").strip()
                        try:
                            h = "#" * int(level)
                        except ValueError:
                            h = "##"
                        parts.append(f"{h} {para.text}")
                    else:
                        parts.append(para.text)
            # Tables
            for table in doc.tables:
                rows = []
                for i, row in enumerate(table.rows):
                    cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows.append("|" + "|".join(["---"] * len(cells)) + "|")
                parts.append("\n".join(rows))
            return "\n\n".join(parts), "docx"
        except Exception as e:
            return f"[docx extraction error: {e}]", "error"

    # ── PPTX ────────────────────────────────────────────────
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(src))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            note = f"pptx-{len(prs.slides)}-slides"
            return "\n\n".join(slides), note
        except Exception as e:
            return f"[pptx extraction error: {e}]", "error"

    # ── XLSX / XLS / XLSM ───────────────────────────────────
    if ext in {".xlsx", ".xlsm"}:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(src), read_only=True, data_only=True)
            sheets = []
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                row_count = 0
                for row in ws.iter_rows(max_row=200, values_only=True):
                    if any(c is not None for c in row):
                        rows.append("\t".join("" if c is None else str(c) for c in row))
                        row_count += 1
                sheets.append(f"## Sheet: {name}\n\n" + "\n".join(rows))
            note = f"xlsx-{len(wb.sheetnames)}-sheets"
            return "\n\n".join(sheets), note
        except Exception as e:
            return f"[xlsx extraction error: {e}]", "error"

    if ext == ".xls":
        try:
            import pandas as pd
            xl = pd.ExcelFile(str(src), engine="xlrd")
            sheets = []
            for name in xl.sheet_names:
                df = xl.parse(name, nrows=200)
                sheets.append(f"## Sheet: {name}\n\n{df.to_markdown(index=False)}")
            note = f"xls-{len(xl.sheet_names)}-sheets"
            return "\n\n".join(sheets), note
        except Exception as e:
            return f"[xls extraction error: {e}]", "error"

    # ── CSV / TSV ────────────────────────────────────────────
    if ext in {".csv", ".tsv"}:
        try:
            import pandas as pd
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(str(src), sep=sep, nrows=200)
            total = sum(1 for _ in open(str(src))) - 1
            note = f"{ext.lstrip('.')}-{total}-rows"
            preview = df.to_markdown(index=False)
            if total > 200:
                preview += f"\n\n*[Showing first 200 of {total} rows]*"
            return preview, note
        except Exception as e:
            return f"[csv extraction error: {e}]", "error"

    # ── PDF ──────────────────────────────────────────────────
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            import base64

            doc = fitz.open(str(src))
            page_count = len(doc)
            pages = []
            scanned_pages = []

            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages.append(("text", i, text))
                else:
                    scanned_pages.append(i)

            # Mixed or fully scanned: rasterize scanned pages as base64 PNG
            # and embed as markdown image-data URIs for multimodal LLM consumption
            if scanned_pages:
                for i, page in enumerate(doc, 1):
                    if i not in scanned_pages:
                        continue
                    mat  = fitz.Matrix(1.5, 1.5)          # 1.5x scale — reasonable quality/size balance
                    pix  = page.get_pixmap(matrix=mat)
                    b64  = base64.b64encode(pix.tobytes("png")).decode()
                    pages.append(("image", i, b64))

            doc.close()

            # Sort all pages by page number and build body
            pages.sort(key=lambda x: x[1])
            body_parts = []
            for kind, num, payload in pages:
                if kind == "text":
                    body_parts.append(f"## Page {num}\n\n{payload}")
                else:
                    # Embedded base64 PNG — multimodal LLMs can read this directly
                    body_parts.append(
                        f"## Page {num}\n\n"
                        f"*[Scanned page — image embedded for multimodal LLM consumption]*\n\n"
                        f"![Page {num}](data:image/png;base64,{payload})"
                    )

            text_count   = sum(1 for k, _, _ in pages if k == "text")
            scanned_count = len(scanned_pages)

            if scanned_count == page_count:
                note = f"pdf-{page_count}p-scanned-images-embedded"
            elif scanned_count > 0:
                note = f"pdf-{page_count}p-mixed-{scanned_count}-scanned"
            else:
                note = f"pdf-{page_count}p-text"

            return "\n\n".join(body_parts), note

        except Exception as e:
            return f"[pdf extraction error: {e}]", "error"

    return f"[unsupported extension: {ext}]", "unknown"


# ─────────────────────────────────────────────────────────────
# LLM enrichment
# ─────────────────────────────────────────────────────────────

def enrich_with_llm(
    filepath: Path,
    content: str,
    type_hint: str,
    endpoint: str,
    model: str,
    api_key: str,
) -> dict:
    truncated = content[:MAX_CONTENT_CHARS]
    if len(content) > MAX_CONTENT_CHARS:
        truncated += "\n... [truncated]"

    prompt = textwrap.dedent(f"""
        You are an OKF enrichment agent. Analyze the file content below and return a JSON
        object with exactly these keys:

        {{
          "type": "<one of: Document, Code Module, Runbook, Configuration, API, Metric, Agent, Dataset, Playbook, Process>",
          "title": "<concise human-readable title, max 8 words>",
          "description": "<one sentence describing what this file contains or does>",
          "tags": ["<tag1>", "<tag2>", "<tag3>"]
        }}

        Rules:
        - Infer type from content, but lean toward the hint: {type_hint}
        - tags: 2-5 lowercase single-word or hyphenated terms reflecting the domain
        - Return ONLY the raw JSON object. No markdown fences, no preamble.

        File: {filepath.name}
        ---
        {truncated}
    """).strip()

    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    payload = {
        "model": model,
        "max_tokens": 300,
        "temperature": 0.1,
        "messages": [{"role": "user", "content": prompt}],
    }

    try:
        resp = requests.post(
            f"{endpoint.rstrip('/')}/chat/completions",
            headers=headers,
            json=payload,
            timeout=30,
        )
        resp.raise_for_status()
        raw = resp.json()["choices"][0]["message"]["content"].strip()
        raw = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE).strip()
        return json.loads(raw)
    except Exception as exc:
        print(f"  [warn] LLM enrichment failed: {exc}")
        return {
            "type": type_hint,
            "title": filepath.stem.replace("_", " ").replace("-", " ").title(),
            "description": f"Content of {filepath.name}",
            "tags": [type_hint.lower().replace(" ", "-")],
        }


# ─────────────────────────────────────────────────────────────
# OKF writers
# ─────────────────────────────────────────────────────────────

def make_frontmatter(meta: dict, extra: dict | None = None) -> str:
    data = {
        "type":        meta.get("type", "Document"),
        "title":       meta.get("title", ""),
        "description": meta.get("description", ""),
        "tags":        meta.get("tags", []),
        "timestamp":   datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
    }
    if extra:
        data.update(extra)
    return "---\n" + yaml.dump(data, default_flow_style=False, allow_unicode=True) + "---\n"


def format_body(content: str, src_ext: str, extraction_note: str) -> str:
    """Wrap extracted content appropriately for OKF body."""
    header = f"# Content\n\n*Source: `{extraction_note}`*\n\n"

    if src_ext in CODE_EXTENSIONS:
        lang = src_ext.lstrip(".")
        return header + f"```{lang}\n{content}\n```\n"

    if src_ext in CONFIG_EXTENSIONS:
        lang = src_ext.lstrip(".") or "text"
        return header + f"```{lang}\n{content}\n```\n"

    # Office/PDF/text: already structured markdown from extractor
    return header + content + "\n"


def write_concept(out_path: Path, frontmatter: str, body: str, dry_run: bool) -> None:
    full = frontmatter + "\n" + body
    if dry_run:
        print(f"  [dry-run] → {out_path}")
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(full, encoding="utf-8")


# ─────────────────────────────────────────────────────────────
# Index / Log
# ─────────────────────────────────────────────────────────────

def write_subdir_index(directory: Path, concepts: list[tuple[Path, dict]], bundle_root: Path, dry_run: bool) -> None:
    lines = [f"# {directory.relative_to(bundle_root)}\n"]
    for p, meta in sorted(concepts, key=lambda x: x[0].name):
        title = meta.get("title", p.stem)
        rel   = p.relative_to(bundle_root)
        lines.append(f"- [{title}](/{rel.as_posix()})")
    index_path = directory / "index.md"
    content = "\n".join(lines) + "\n"
    if dry_run:
        print(f"  [dry-run] → {index_path}")
    else:
        index_path.write_text(content, encoding="utf-8")


def write_root_index(bundle_root: Path, all_concepts: list[tuple[Path, dict]], source_dir: Path, dry_run: bool) -> None:
    ts = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    lines = [
        "---\nokf_version: \"0.1\"\n---\n",
        "# Knowledge Bundle\n",
        f"Generated from `{source_dir.resolve()}` on {ts}.\n",
        f"**Concepts:** {len(all_concepts)}\n",
        "\n## Contents\n",
    ]

    by_dir: dict[str, list] = {}
    for path, meta in all_concepts:
        rel   = path.relative_to(bundle_root)
        group = rel.parts[0] if len(rel.parts) > 1 else "."
        by_dir.setdefault(group, []).append((path, meta))

    for group, items in sorted(by_dir.items()):
        label = "Root" if group == "." else f"{group}/"
        lines.append(f"### {label}\n")
        for p, m in sorted(items, key=lambda x: x[0].name):
            rel = p.relative_to(bundle_root)
            lines.append(f"- [{m.get('title', p.stem)}](/{rel.as_posix()})")
        lines.append("")

    content = "\n".join(lines)
    index_path = bundle_root / "index.md"
    if dry_run:
        print(f"  [dry-run] → {index_path}")
    else:
        index_path.write_text(content, encoding="utf-8")


def write_log(bundle_root: Path, converted: list[tuple[Path, dict]], dry_run: bool) -> None:
    date = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    lines = [
        "# Change Log\n",
        f"## {date}\n",
        f"**Initialization**: Generated OKF bundle. {len(converted)} concept(s) created.\n",
    ]
    for p, meta in sorted(converted, key=lambda x: x[0].name):
        rel = p.relative_to(bundle_root)
        lines.append(f"- **Creation**: [{meta.get('title', p.stem)}](/{rel.as_posix()})")
    log_path = bundle_root / "log.md"
    content = "\n".join(lines) + "\n"
    if dry_run:
        print(f"  [dry-run] → {log_path}")
    else:
        log_path.write_text(content, encoding="utf-8")


# ─────────────────────────────────────────────────────────────
# ─────────────────────────────────────────────────────────────
# Manifest — tracks which source files have been processed
# ─────────────────────────────────────────────────────────────

def file_hash(path: Path) -> str:
    """SHA-256 of file contents — detects changes."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def load_manifest(bundle_root: Path) -> dict:
    mp = bundle_root / MANIFEST_FILE
    if mp.exists():
        try:
            return json.loads(mp.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def save_manifest(bundle_root: Path, manifest: dict, dry_run: bool) -> None:
    mp = bundle_root / MANIFEST_FILE
    if dry_run:
        print(f"  [dry-run] → {mp}")
    else:
        mp.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


# ─────────────────────────────────────────────────────────────
# File discovery
# ─────────────────────────────────────────────────────────────

def collect_files(source_dir: Path) -> list[Path]:
    files = []
    for root, dirs, filenames in os.walk(source_dir):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fname in filenames:
            p = Path(root) / fname
            if p.suffix.lower() in ALL_SUPPORTED:
                files.append(p)
    return sorted(files)


def detect_changes(
    source_dir: Path,
    manifest: dict,
) -> tuple[list[Path], list[Path], list[str]]:
    """
    Returns:
        new_files     — not in manifest
        changed_files — hash differs from manifest
        deleted_keys  — in manifest but source no longer exists
    """
    all_src  = collect_files(source_dir)
    src_keys = {str(p.relative_to(source_dir)): p for p in all_src}

    new_files:    list[Path] = []
    changed_files: list[Path] = []
    deleted_keys:  list[str]  = []

    for key, src_path in src_keys.items():
        current_hash = file_hash(src_path)
        if key not in manifest:
            new_files.append(src_path)
        elif manifest[key]["hash"] != current_hash:
            changed_files.append(src_path)

    for key in manifest:
        if key not in src_keys:
            deleted_keys.append(key)

    return new_files, changed_files, deleted_keys


# ─────────────────────────────────────────────────────────────
# Incremental log append
# ─────────────────────────────────────────────────────────────

def append_log(
    bundle_root: Path,
    new_concepts:     list[tuple[Path, dict]],
    updated_concepts: list[tuple[Path, dict]],
    deleted_keys:     list[str],
    dry_run: bool,
) -> None:
    """Prepend a new dated sync section into log.md, keeping all history."""
    log_path = bundle_root / "log.md"
    existing = log_path.read_text(encoding="utf-8") if log_path.exists() else "# Change Log\n"

    date = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    ts   = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    lines = [f"\n## {date} — sync at {ts}\n"]

    for p, meta in sorted(new_concepts, key=lambda x: x[0].name):
        rel = p.relative_to(bundle_root)
        lines.append(f"- **Creation**: [{meta.get('title', p.stem)}](/{rel.as_posix()})")

    for p, meta in sorted(updated_concepts, key=lambda x: x[0].name):
        rel = p.relative_to(bundle_root)
        lines.append(f"- **Update**: [{meta.get('title', p.stem)}](/{rel.as_posix()})")

    for key in sorted(deleted_keys):
        lines.append(f"- **Deletion**: `{key}` removed from source")

    if not (new_concepts or updated_concepts or deleted_keys):
        lines.append("- No changes detected.")

    # Insert new section right after the "# Change Log" header line
    header_end = existing.find("\n", existing.find("# Change Log")) + 1
    new_content = existing[:header_end] + "\n".join(lines) + "\n\n" + existing[header_end:].lstrip("\n")

    if dry_run:
        print(f"  [dry-run] would prepend to log.md:")
        for line in lines:
            print(f"    {line}")
    else:
        log_path.write_text(new_content, encoding="utf-8")


def rebuild_root_index(bundle_root: Path, source_dir: Path, manifest: dict, dry_run: bool) -> None:
    """Rebuild root index.md from all concept files currently in the bundle."""
    # Collect all .md files that are concept files (not reserved)
    all_concepts = []
    for md_path in sorted(bundle_root.rglob("*.md")):
        rel = md_path.relative_to(bundle_root)
        # Skip reserved files and hidden files
        if rel.parts[0].startswith("."):
            continue
        if md_path.stem.lower() in RESERVED_FILENAMES:
            continue
        # Read frontmatter to get title/type
        try:
            raw = md_path.read_text(encoding="utf-8")
            if raw.startswith("---"):
                fm_end = raw.find("---", 3)
                fm = yaml.safe_load(raw[3:fm_end])
                meta = {"title": fm.get("title", md_path.stem), "type": fm.get("type", "")}
            else:
                meta = {"title": md_path.stem, "type": "Document"}
        except Exception:
            meta = {"title": md_path.stem, "type": "Document"}
        all_concepts.append((md_path, meta))

    write_root_index(bundle_root, all_concepts, source_dir, dry_run)

    # Also rebuild per-subdirectory indexes
    by_subdir: dict[Path, list] = {}
    for p, meta in all_concepts:
        if p.parent != bundle_root:
            by_subdir.setdefault(p.parent, []).append((p, meta))
    for subdir, concepts in by_subdir.items():
        write_subdir_index(subdir, concepts, bundle_root, dry_run)


# ─────────────────────────────────────────────────────────────
# Process a single file (shared by convert and sync)
# ─────────────────────────────────────────────────────────────

def process_file(
    src_path: Path,
    source_dir: Path,
    output_dir: Path,
    endpoint: str,
    model: str,
    api_key: str,
    dry_run: bool,
) -> tuple[Path, dict]:
    rel      = src_path.relative_to(source_dir)
    out_stem = rel.with_suffix(".md")
    out_path = output_dir / out_stem

    if out_path.stem.lower() in RESERVED_FILENAMES:
        out_path = output_dir / rel.parent / f"_{rel.stem}.md"

    ext      = src_path.suffix.lower()
    category = "text" if ext in TEXT_EXTENSIONS else ("office" if ext in OFFICE_EXTENSIONS else "pdf")
    print(f"  ({category}) {rel}")

    content, extraction_note = extract_text(src_path)
    print(f"  Extracted {len(content):,} chars  [{extraction_note}]")

    type_hint = EXTENSION_TYPE_HINT.get(ext, "Document")
    print(f"  Enriching with LLM ({model})...")
    meta = enrich_with_llm(src_path, content, type_hint, endpoint, model, api_key)
    print(f"  type={meta['type']} | title={meta['title']}")

    fm   = make_frontmatter(meta, extra={"source_file": str(rel), "source_format": ext.lstrip(".")})
    body = format_body(content, ext, extraction_note)
    write_concept(out_path, fm, body, dry_run)
    print()

    return out_path, meta


# ─────────────────────────────────────────────────────────────
# Full convert (initial run)
# ─────────────────────────────────────────────────────────────

def convert(source_dir: Path, output_dir: Path, endpoint: str, model: str, api_key: str, dry_run: bool) -> None:
    files = collect_files(source_dir)
    if not files:
        print("No supported files found in input directory.")
        sys.exit(1)

    text_count   = sum(1 for f in files if f.suffix.lower() in TEXT_EXTENSIONS)
    office_count = sum(1 for f in files if f.suffix.lower() in OFFICE_EXTENSIONS)
    pdf_count    = sum(1 for f in files if f.suffix.lower() in PDF_EXTENSIONS)
    print(f"\nFound {len(files)} file(s): {text_count} text, {office_count} office, {pdf_count} PDF\n")

    if not dry_run:
        output_dir.mkdir(parents=True, exist_ok=True)

    all_concepts: list[tuple[Path, dict]] = []
    manifest: dict = {}

    for i, src_path in enumerate(files, 1):
        print(f"[{i}/{len(files)}]")
        out_path, meta = process_file(src_path, source_dir, output_dir, endpoint, model, api_key, dry_run)
        all_concepts.append((out_path, meta))
        key = str(src_path.relative_to(source_dir))
        manifest[key] = {"hash": file_hash(src_path), "concept": str(out_path.relative_to(output_dir))}

    # Write root + subdir indexes
    by_subdir: dict[Path, list] = {}
    for p, meta in all_concepts:
        if p.parent != output_dir:
            by_subdir.setdefault(p.parent, []).append((p, meta))
    for subdir, concepts in by_subdir.items():
        write_subdir_index(subdir, concepts, output_dir, dry_run)
    write_root_index(output_dir, all_concepts, source_dir, dry_run)
    write_log(output_dir, all_concepts, dry_run)
    save_manifest(output_dir, manifest, dry_run)

    print(f"{'[dry-run] ' if dry_run else ''}Done!")
    print(f"  OKF bundle → {output_dir.resolve()}")
    print(f"  {len(all_concepts)} concept(s) | index.md | log.md | {MANIFEST_FILE}")


# ─────────────────────────────────────────────────────────────
# Incremental sync
# ─────────────────────────────────────────────────────────────

def sync(source_dir: Path, output_dir: Path, endpoint: str, model: str, api_key: str, dry_run: bool) -> None:
    if not output_dir.exists():
        print("Bundle directory does not exist. Run without --sync to do the initial conversion first.")
        sys.exit(1)

    manifest = load_manifest(output_dir)
    if not manifest:
        print("No manifest found in bundle. Run without --sync to do the initial conversion first.")
        sys.exit(1)

    new_files, changed_files, deleted_keys = detect_changes(source_dir, manifest)

    print(f"\nSync status:")
    print(f"  New     : {len(new_files)} file(s)")
    print(f"  Changed : {len(changed_files)} file(s)")
    print(f"  Deleted : {len(deleted_keys)} file(s)")

    if not new_files and not changed_files and not deleted_keys:
        print("\nNothing to do — bundle is up to date.")
        return

    to_process = new_files + changed_files
    new_concepts:     list[tuple[Path, dict]] = []
    updated_concepts: list[tuple[Path, dict]] = []

    if to_process:
        print()
    for i, src_path in enumerate(to_process, 1):
        is_new = src_path in new_files
        label  = "NEW" if is_new else "CHANGED"
        print(f"[{i}/{len(to_process)}] [{label}]")
        out_path, meta = process_file(src_path, source_dir, output_dir, endpoint, model, api_key, dry_run)

        key = str(src_path.relative_to(source_dir))
        manifest[key] = {"hash": file_hash(src_path), "concept": str(out_path.relative_to(output_dir))}

        if is_new:
            new_concepts.append((out_path, meta))
        else:
            updated_concepts.append((out_path, meta))

    # Handle deletions: remove concept file and manifest entry
    for key in deleted_keys:
        concept_rel = manifest[key].get("concept", "")
        concept_path = output_dir / concept_rel if concept_rel else None
        if concept_path and concept_path.exists():
            if dry_run:
                print(f"  [dry-run] would delete concept → {concept_path}")
            else:
                concept_path.unlink()
                print(f"  Deleted concept: {concept_path.relative_to(output_dir)}")
        del manifest[key]

    # Rebuild indexes and append to log
    rebuild_root_index(output_dir, source_dir, manifest, dry_run)
    append_log(output_dir, new_concepts, updated_concepts, deleted_keys, dry_run)
    save_manifest(output_dir, manifest, dry_run)

    total_changes = len(new_concepts) + len(updated_concepts) + len(deleted_keys)
    print(f"\n{'[dry-run] ' if dry_run else ''}Sync complete — {total_changes} change(s) applied.")
    print(f"  +{len(new_concepts)} created  ~{len(updated_concepts)} updated  -{len(deleted_keys)} deleted")


# ─────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Convert or sync a directory into an OKF-conformant knowledge bundle.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""
            Supported formats:
              Text   : .md .txt .rst .py .js .ts .sh .yaml .yml .json .toml .cfg .ini .env
              Office : .docx .pptx .xlsx .xls .xlsm .csv .tsv
              PDF    : .pdf (text-layer + scanned pages embedded as base64 images)

            Modes:
              Initial conversion (first run):
                python okf_convert.py --input ./docs --output ./okf-bundle

              Incremental sync (after adding/changing/deleting files):
                python okf_convert.py --input ./docs --output ./okf-bundle --sync

              Preview without writing:
                python okf_convert.py --input ./docs --output ./okf-bundle --sync --dry-run

            The --sync flag:
              - Detects new files  → converts and adds to index + log (Creation entry)
              - Detects changed files → re-converts and updates log (Update entry)
              - Detects deleted files → removes concept and logs deletion (Deletion entry)
              - Leaves unchanged files untouched
              - Stores hashes in .okf_manifest.json inside the bundle
        """),
    )
    parser.add_argument("--input",    required=True,             help="Source directory")
    parser.add_argument("--output",   required=True,             help="OKF bundle directory")
    parser.add_argument("--endpoint", default=DEFAULT_ENDPOINT,  help=f"LiteLLM base URL (default: {DEFAULT_ENDPOINT})")
    parser.add_argument("--model",    default=DEFAULT_MODEL,     help=f"Model for enrichment (default: {DEFAULT_MODEL})")
    parser.add_argument("--api-key",  default=os.environ.get("LITELLM_API_KEY", ""), help="API key or set LITELLM_API_KEY")
    parser.add_argument("--sync",     action="store_true",       help="Incremental sync: only process new/changed files")
    parser.add_argument("--dry-run",  action="store_true",       help="Print what would happen without writing files")

    args   = parser.parse_args()
    source = Path(args.input)
    output = Path(args.output)

    if not source.exists() or not source.is_dir():
        print(f"Error: input directory not found: {source}")
        sys.exit(1)

    mode = "sync" if args.sync else "convert"
    print(f"\nOKF Converter  [mode: {mode}]")
    print(f"  Input    : {source.resolve()}")
    print(f"  Output   : {output.resolve()}")
    print(f"  Endpoint : {args.endpoint}")
    print(f"  Model    : {args.model}")
    print(f"  Dry run  : {args.dry_run}")

    if args.sync:
        sync(
            source_dir=source,
            output_dir=output,
            endpoint=args.endpoint,
            model=args.model,
            api_key=args.api_key,
            dry_run=args.dry_run,
        )
    else:
        if not args.dry_run and output.exists() and any(output.iterdir()):
            ans = input(f"\nOutput directory not empty: {output}\nContinue and overwrite? [y/N] ").strip().lower()
            if ans != "y":
                sys.exit(0)
        convert(
            source_dir=source,
            output_dir=output,
            endpoint=args.endpoint,
            model=args.model,
            api_key=args.api_key,
            dry_run=args.dry_run,
        )


if __name__ == "__main__":
    main()
